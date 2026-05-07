#!/usr/bin/env python3
"""Generate the architecture / pipeline overview figure for the paper."""

from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Rectangle
from matplotlib.lines import Line2D


PAGE_W = 14.0   # inches
PAGE_H = 6.5

COL_INPUT  = "#DCE6F1"
COL_ASR    = "#EAEAEA"
COL_PROMPT = "#E2EFDA"
COL_MODEL  = "#FCE4D6"
COL_HEAD   = "#FFF2CC"
COL_ABLATE = "#FFFAE0"
EDGE       = "#3B3B3B"
TEXT       = "#1A1A1A"
ARROW      = "#555555"

ICML_RC = {
    "font.family":        "serif",
    "font.serif":         ["Times New Roman", "DejaVu Serif"],
    "mathtext.fontset":   "dejavuserif",
    "font.size":          11,
    "axes.titlesize":     12,
    "savefig.bbox":       "tight",
    "savefig.pad_inches": 0.05,
    "savefig.dpi":        300,
}


def rounded_box(ax, x, y, w, h, color, edge=EDGE, lw=1.0, zorder=2):
    box = FancyBboxPatch((x, y), w, h,
                         boxstyle="round,pad=0.02,rounding_size=0.10",
                         linewidth=lw, edgecolor=edge, facecolor=color,
                         zorder=zorder)
    ax.add_patch(box)
    return box


def title_text(ax, x, y, text, **kw):
    kw.setdefault("color", TEXT)
    kw.setdefault("weight", "bold")
    kw.setdefault("ha", "center")
    kw.setdefault("va", "bottom")
    kw.setdefault("fontsize", 10.5)
    ax.text(x, y, text, **kw)


def body_text(ax, x, y, text, fontsize=8.8, ha="center", va="center", **kw):
    kw.setdefault("color", TEXT)
    ax.text(x, y, text, ha=ha, va=va, fontsize=fontsize, **kw)


def arrow(ax, x0, y0, x1, y1, lw=1.4, style="-|>", mutation=15,
          color=ARROW, ls="-"):
    ax.add_patch(FancyArrowPatch((x0, y0), (x1, y1),
                                  arrowstyle=style,
                                  mutation_scale=mutation,
                                  linewidth=lw, color=color,
                                  linestyle=ls,
                                  zorder=3))


def draw_architecture(ax):
    ax.set_xlim(0, PAGE_W)
    ax.set_ylim(0, PAGE_H)
    ax.set_aspect("equal")
    ax.axis("off")

    Y_MAIN = 3.6           # vertical center for the main flow
    BH     = 1.6           # main-stage box height

    x = 0.25
    rounded_box(ax, x, 4.7, 1.9, 1.0, COL_INPUT)
    title_text(ax, x + 0.95, 5.5, "CTAF Audio")
    body_text(ax, x + 0.95, 5.05, "MP3, 30–90 s\nradio chatter", fontsize=8.0)

    rounded_box(ax, x, 3.3, 1.9, 1.0, COL_INPUT)
    title_text(ax, x + 0.95, 4.10, "METAR")
    body_text(ax, x + 0.95, 3.65, "raw + parsed\nweather report", fontsize=8.0)

    # Section header
    body_text(ax, x + 0.95, 6.05, "(a) Inputs", fontsize=9.2, weight="bold",
               va="bottom")

    x = 2.55
    rounded_box(ax, x, 4.7, 1.85, 1.0, COL_ASR)
    title_text(ax, x + 0.92, 5.5, "Whisper ASR")
    body_text(ax, x + 0.92, 5.05, "large-v3 default\n(base/medium/large)",
               fontsize=8.0)
    body_text(ax, x + 0.92, 6.05, "(b) Speech → text", fontsize=9.2,
               weight="bold", va="bottom")

    arrow(ax, 2.15, 5.20, 2.55, 5.20)        # audio → ASR
    # METAR bypass (curve straight to prompt)
    arrow(ax, 2.15, 3.80, 4.80, 3.80, style="-|>")

    # Output of ASR feeds prompt
    arrow(ax, 4.40, 5.20, 4.80, 4.50)         # ASR → prompt assembly

    x = 4.80
    rounded_box(ax, x, 3.05, 2.55, 2.65, COL_PROMPT)
    title_text(ax, x + 1.275, 5.5, "Prompt Assembly")
    body_text(ax, x + 1.275, 5.10,
               "• System prompt\n"
               "• ICL exemplars: 0 / 1 / 2 per class\n"
               "• Optional CoT cue\n"
               "• Concatenate METAR + transcript",
               fontsize=8.4, ha="center", va="top")
    body_text(ax, x + 1.275, 6.05, "(c) Conditioning", fontsize=9.2,
               weight="bold", va="bottom")

    # Two-task switch as a clean inline badge with text rendered last
    body_text(ax, x + 1.275, 3.42,
               "Task switch:  3-class  |  Binary",
               fontsize=8.6, weight="bold", color="#222")

    arrow(ax, 7.35, 4.40, 7.85, 4.40)
    x = 7.85
    rounded_box(ax, x, 3.05, 2.85, 2.65, COL_MODEL)
    title_text(ax, x + 1.425, 5.5, "Safety Classifier (LLM)")
    body_text(ax, x + 0.20, 5.10,
               "Open-source\n"
               "• Qwen 2.5-7B-Instruct\n"
               "• Mistral-7B-Instruct\n"
               "• Gemma-2-9B-IT",
               fontsize=8.2, ha="left", va="top")
    body_text(ax, x + 1.55, 5.10,
               "Closed-source\n"
               "• GPT-4o\n"
               "• GPT-5.4\n"
               "• Claude Sonnet 4.6",
               fontsize=8.2, ha="left", va="top")

    body_text(ax, x + 1.425, 3.30,
               r"output: $\{$label, conf, reasoning, $p_{\mathrm{class}}$$\}$",
               fontsize=8.2, weight="bold")
    body_text(ax, x + 1.425, 6.05, "(d) Inference", fontsize=9.2,
               weight="bold", va="bottom")

    arrow(ax, 10.70, 4.40, 11.20, 4.40)
    x = 11.20
    rounded_box(ax, x, 3.05, 2.55, 2.65, COL_HEAD)
    title_text(ax, x + 1.275, 5.5, "Evaluation Heads")
    body_text(ax, x + 1.275, 5.05,
               "3-class:  Nominal | Warning | Hazard\n"
               "Binary:   Nominal | Danger\n",
               fontsize=8.4)
    body_text(ax, x + 1.275, 4.20,
               "Metrics:  Accuracy, Macro-F$_1$,\n"
               "per-class P/R/F$_1$, Confusion Mat.",
               fontsize=8.4)
    body_text(ax, x + 1.275, 3.40,
               "Binary only:  AUROC, AUPRC,\n"
               "PR + ROC curves",
               fontsize=8.4, style="italic")
    body_text(ax, x + 1.275, 6.05, "(e) Scoring", fontsize=9.2,
               weight="bold", va="bottom")

    rounded_box(ax, 0.25, 0.55, 13.50, 1.50, COL_ABLATE, lw=1.0)
    body_text(ax, 0.45, 1.85, "(f) Robustness ablations",
               fontsize=10.5, weight="bold", ha="left", va="center")

    body_text(ax, 1.25, 1.30,
               "ASR quality\n"
               "Whisper {base, medium, large-v3}",
               fontsize=8.6, ha="left", va="top")

    body_text(ax, 5.65, 1.30,
               "Audio noise\n"
               r"NSR $\in$ {5, 10, 25, 50, 75}%" "\n"
               "(re-transcribe, then classify)",
               fontsize=8.6, ha="left", va="top")

    body_text(ax, 9.95, 1.30,
               "Text masking\n"
               r"rate $\in$ {10, 20, 40, 60, 80}% $\times$ {word, utterance}",
               fontsize=8.6, ha="left", va="top")

    for x_from, x_to, label_x in ((2.20, 3.45, 1.30),
                                   (5.85, 1.20, 5.85),
                                   (10.65, 6.00, 10.65)):
        arrow(ax, label_x, 2.05, x_to, 3.05, lw=0.8, style="-|>",
              mutation=10, ls="--", color="#888")

    # Caption sits below the ablation band, away from section headers.
    body_text(ax, PAGE_W / 2, 0.30,
               "Per scenario:  one (METAR, audio, ground-truth label) tuple "
               "from the CTAF-KHAF benchmark — 94 test scenarios + 6 ICL.",
               fontsize=8.4, style="italic", color="#555")


def render_pdf_and_png(out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)
    with mpl.rc_context(ICML_RC):
        fig, ax = plt.subplots(figsize=(PAGE_W, PAGE_H))
        draw_architecture(ax)
        fig.savefig(out_dir / "architecture.pdf")
        fig.savefig(out_dir / "architecture.png", dpi=300)
        plt.close(fig)
        print(f"  Saved → {out_dir / 'architecture.pdf'}")
        print(f"  Saved → {out_dir / 'architecture.png'}")


def render_pptx(out_dir: Path):
    """Build a one-slide editable PowerPoint with the same layout."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  python-pptx not installed; skipping .pptx export.")
        print("    pip install python-pptx")
        return

    def hex2rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = Presentation()
    prs.slide_width  = Inches(14)
    prs.slide_height = Inches(6.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    def add_box(x, y, w, h, fill, text, bold_first_line=True, fontsize=9):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x), Inches(y),
                                        Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex2rgb(fill)
        shape.line.color.rgb = hex2rgb(EDGE)
        shape.line.width = Pt(0.75)
        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = Pt(fontsize + (1.5 if (bold_first_line and i == 0) else 0))
            r.font.bold = bold_first_line and i == 0
            r.font.name = "Times New Roman"
            r.font.color.rgb = hex2rgb(TEXT)
        return shape

    def add_label(x, y, w, h, text, fontsize=11, bold=False, italic=False,
                  align=PP_ALIGN.LEFT, color=TEXT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y),
                                       Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(fontsize)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = "Times New Roman"
        r.font.color.rgb = hex2rgb(color)
        return tb

    def add_arrow(x1, y1, x2, y2, dashed=False):
        line = slide.shapes.add_connector(2,  # straight connector
                                           Inches(x1), Inches(y1),
                                           Inches(x2), Inches(y2))
        line.line.color.rgb = hex2rgb(ARROW)
        line.line.width = Pt(1.25)
        if dashed:
            try:
                from pptx.oxml.ns import qn
                ln = line.line._get_or_add_ln()
                ln.set("cap", "flat")
                lst = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
                ln.append(lst)
            except Exception:
                pass
        # arrowhead
        ln = line.line
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            ln_xml = ln._get_or_add_ln()
            tail = etree.SubElement(ln_xml, qn("a:tailEnd"),
                                     {"type": "triangle", "w": "med", "len": "med"})
        except Exception:
            pass
        return line

    # Caption strip at top
    add_label(0.3, 0.10, 13.5, 0.30,
              "Per scenario: one (METAR, audio, ground-truth label) tuple "
              "from the CTAF-KHAF benchmark (94 test scenarios + 6 ICL).",
              fontsize=9, italic=True, color="#444444", align=PP_ALIGN.CENTER)

    # Main flow row
    add_label(0.25, 0.45, 1.9, 0.25, "(a) Inputs", fontsize=10, bold=True,
              align=PP_ALIGN.CENTER)
    add_box(0.25, 0.75, 1.9, 1.0, COL_INPUT, "CTAF Audio\nMP3, 30–90 s")
    add_box(0.25, 2.10, 1.9, 1.0, COL_INPUT, "METAR\nraw + parsed weather")

    add_label(2.55, 0.45, 1.85, 0.25, "(b) ASR", fontsize=10, bold=True,
              align=PP_ALIGN.CENTER)
    add_box(2.55, 0.75, 1.85, 1.0, COL_ASR,
             "Whisper ASR\nlarge-v3 default")

    add_label(4.80, 0.45, 2.55, 0.25, "(c) Conditioning", fontsize=10, bold=True,
              align=PP_ALIGN.CENTER)
    add_box(4.80, 0.75, 2.55, 2.65, COL_PROMPT,
             "Prompt Assembly\n"
             "• System prompt\n"
             "• ICL exemplars (0/1/2 per class)\n"
             "• Optional CoT cue\n"
             "• Concatenate METAR + transcript\n\n"
             "Task switch: 3-class | Binary",
             fontsize=9)

    add_label(7.85, 0.45, 2.85, 0.25, "(d) Inference", fontsize=10, bold=True,
              align=PP_ALIGN.CENTER)
    add_box(7.85, 0.75, 2.85, 2.65, COL_MODEL,
             "Safety Classifier (LLM)\n"
             "\n"
             "Open-source:\n"
             "  Qwen 2.5-7B,  Mistral-7B,  Gemma-2-9B\n"
             "\n"
             "Closed-source:\n"
             "  GPT-4o,  GPT-5.4,  Claude Sonnet 4.6\n"
             "\n"
             "output: {label, conf, reasoning, p_class}",
             fontsize=9)

    add_label(11.20, 0.45, 2.55, 0.25, "(e) Scoring", fontsize=10, bold=True,
              align=PP_ALIGN.CENTER)
    add_box(11.20, 0.75, 2.55, 2.65, COL_HEAD,
             "Evaluation Heads\n"
             "\n"
             "3-class: Nominal | Warning | Hazard\n"
             "Binary:  Nominal | Danger\n"
             "\n"
             "Metrics:  Acc, Macro-F1, P/R/F1, CM\n"
             "Binary only:  AUROC, AUPRC,\n"
             "PR + ROC curves",
             fontsize=9)

    # arrows between main blocks
    add_arrow(2.15, 1.25, 2.55, 1.25)              # audio → ASR
    add_arrow(2.15, 2.60, 4.80, 2.60)              # METAR → prompt
    add_arrow(4.40, 1.25, 4.80, 1.85)              # ASR → prompt
    add_arrow(7.35, 2.10, 7.85, 2.10)              # prompt → model
    add_arrow(10.70, 2.10, 11.20, 2.10)            # model → heads

    # Ablations band
    add_box(0.25, 4.45, 13.50, 1.50, COL_ABLATE,
             "(f) Robustness ablations", bold_first_line=False)
    # Overlay the bold heading nicer:
    add_label(0.45, 4.50, 4.0, 0.30, "(f) Robustness ablations",
              fontsize=11, bold=True)
    add_label(0.45, 4.95, 4.0, 0.95,
              "ASR quality\nWhisper {base, medium, large-v3}",
              fontsize=9)
    add_label(4.85, 4.95, 4.5, 0.95,
              "Audio noise\nNSR ∈ {5, 10, 25, 50, 75}%\n"
              "(re-transcribe, then classify)",
              fontsize=9)
    add_label(9.80, 4.95, 4.0, 0.95,
              "Text masking\nrate ∈ {10, 20, 40, 60, 80}% × {word, utterance}",
              fontsize=9)

    # Dashed arrows from ablations up to the input/ASR side
    for x in (1.20, 5.60, 9.90):
        add_arrow(x, 4.45, x, 3.55, dashed=True)

    out_path = out_dir / "architecture.pptx"
    prs.save(out_path)
    print(f"  Saved → {out_path}")


def main():
    out_dir = Path("figures")
    print("Generating architecture figure...")
    render_pdf_and_png(out_dir)
    print("Generating PowerPoint...")
    render_pptx(out_dir)
    print("\nDone. Files in figures/:")
    for f in sorted(out_dir.glob("architecture.*")):
        print(f"  {f}")


if __name__ == "__main__":
    main()
